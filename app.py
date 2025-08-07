import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import io
from collections import defaultdict

# --- Configuração da Página ---
st.set_page_config(layout="wide", page_title="Gerador de Organograma Editável")
st.title("Gerador de Organograma Editável para PowerPoint 🏢")

# --- Assinatura ---
st.markdown("<p style='font-size:16px; color:gray;'>by <strong>Geaco - Com &amp; Serv</strong></p>", unsafe_allow_html=True)

st.info("Este aplicativo gera um organograma com formas e conectores que podem ser editados diretamente no PowerPoint.")

# --- Inicialização do Estado da Aplicação ---
if 'relationships' not in st.session_state:
    st.session_state.relationships = []

col1, col2 = st.columns([1, 1.2])
with col1:
    st.header("1. Adicione as Relações Societárias")
    with st.form("relationship_form", clear_on_submit=True):
        controladora = st.text_input("Nome da Empresa Controladora/Holding")
        subsidiaria = st.text_input("Nome da Empresa Subsidiária/Afiliada")
        percentual = st.number_input("Percentual de Posse (%)", min_value=0.0, max_value=100.0, step=0.1, format="%.1f")
        
        submitted = st.form_submit_button("➕ Adicionar Relação")
        if submitted:
            if not controladora or not subsidiaria:
                st.warning("Por favor, preencha o nome da controladora e da subsidiária.")
            else:
                st.session_state.relationships.append({
                    "Controladora": controladora.strip(),
                    "Subsidiária": subsidiaria.strip(),
                    "Percentual": percentual
                })
                st.success(f"Adicionado: {controladora} → {percentual}% de {subsidiaria}")

with col2:
    st.header("2. Estrutura Atual")
    if st.session_state.relationships:
        header_cols = st.columns([3, 3, 2, 1])
        header_cols[0].markdown("**Controladora**")
        header_cols[1].markdown("**Subsidiária**")
        header_cols[2].markdown("**Percentual**")
        header_cols[3].markdown("**Ação**")
        st.markdown("---")

        for index, rel in enumerate(st.session_state.relationships):
            row_cols = st.columns([3, 3, 2, 1])
            row_cols[0].write(rel["Controladora"])
            row_cols[1].write(rel["Subsidiária"])
            row_cols[2].write(f"{rel['Percentual']}%")
            
            if row_cols[3].button("🗑️", key=f"delete_{index}", help="Excluir esta relação"):
                st.session_state.relationships.pop(index)
                st.rerun()

        if st.button("🧹 Limpar Tudo", use_container_width=True):
            st.session_state.relationships = []
            st.rerun()
    else:
        st.info("Nenhuma relação adicionada ainda.")

st.markdown("---")
st.header("3. Gerar a Apresentação Editável")

def build_tree(relationships):
    tree = defaultdict(lambda: {'children': [], 'data': None})
    all_nodes, children_nodes = set(), set()
    for rel in relationships:
        parent, child, percent = rel['Controladora'], rel['Subsidiária'], rel['Percentual']
        tree[parent]['children'].append({'name': child, 'percent': percent})
        all_nodes.update([parent, child])
        children_nodes.add(child)
    root_nodes = list(all_nodes - children_nodes)
    if not root_nodes and all_nodes:
        root_nodes = [list(all_nodes)[0]]
    return tree, root_nodes

def calculate_positions_recursive(node_name, tree, level, positions, x_offset):
    BOX_WIDTH, BOX_HEIGHT = Inches(2.5), Inches(1.2)
    H_SPACING, V_SPACING = Inches(0.5), Inches(1.5)
    children = tree[node_name]['children']
    total_width = len(children) * BOX_WIDTH + max(0, len(children) - 1) * H_SPACING
    x = x_offset
    y = level * (BOX_HEIGHT + V_SPACING)
    positions[node_name] = {'x': x + total_width / 2 - BOX_WIDTH / 2, 'y': y, 'width': BOX_WIDTH, 'height': BOX_HEIGHT}
    child_x = x
    for child_info in children:
        child_name = child_info['name']
        calculate_positions_recursive(child_name, tree, level + 1, positions, child_x)
        child_x += BOX_WIDTH + H_SPACING

def draw_organogram(slide, relationships, positions, tree):
    shapes = {}
    for name, pos in positions.items():
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(pos['x']), int(pos['y']), int(pos['width']), int(pos['height']))
        SHAPE_COLOR = RGBColor(250, 190, 80)
        shape.fill.solid()
        shape.fill.fore_color.rgb = SHAPE_COLOR
        line = shape.line
        line.color.rgb = SHAPE_COLOR
        line.width = Pt(1.5)
        shape.text = name
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.name = 'Aptos Display'
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.font.shadow = True
        shapes[name] = shape

    for rel in relationships:
        parent_name, child_name, percent = rel['Controladora'], rel['Subsidiária'], rel['Percentual']
        if parent_name in shapes and child_name in shapes:
            from_shape, to_shape = shapes[parent_name], shapes[child_name]
            connector = slide.shapes.add_connector(
                MSO_CONNECTOR.ELBOW,
                int(from_shape.left + from_shape.width / 2),
                int(from_shape.top + from_shape.height),
                int(to_shape.left + to_shape.width / 2),
                int(to_shape.top)
            )
            connector.begin_connect(from_shape, 2)
            connector.end_connect(to_shape, 0)
            connector.line.color.rgb = RGBColor(0, 0, 0)
            connector.line.width = Pt(1.5)

            mid_x = int((from_shape.left + to_shape.left) / 2)
            mid_y = int((from_shape.top + to_shape.top) / 2)
            textbox = slide.shapes.add_textbox(mid_x - int(Inches(0.2)), mid_y - int(Inches(0.1)), int(Inches(0.4)), int(Inches(0.2)))
            textbox.text = f"{percent}%"
            p = textbox.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(0, 0, 0)
            textbox.fill.background()
            textbox.line.fill.background()

if st.session_state.relationships:
    if st.button("🚀 Gerar Apresentação Editável", type="primary"):
        with st.spinner("Construindo organograma editável... Isso pode levar um momento."):
            tree, root_nodes = build_tree(st.session_state.relationships)
            positions = {}
            current_x_offset = Inches(0.5)
            for root_name in root_nodes:
                calculate_positions_recursive(root_name, tree, 0, positions, current_x_offset)
                max_x = max(pos['x'] + pos['width'] for pos in positions.values())
                current_x_offset = max_x + Inches(1.0)

            SLIDE_WIDTH = Inches(10)
            min_x = min(pos['x'] for pos in positions.values())
            max_x = max(pos['x'] + pos['width'] for pos in positions.values())
            total_width = max_x - min_x

            if total_width < SLIDE_WIDTH:
                shift_x = (SLIDE_WIDTH - total_width) / 2 - min_x
                for pos in positions.values():
                    pos['x'] += shift_x

            prs = Presentation()
            prs.slide_width = SLIDE_WIDTH
            prs.slide_height = Inches(7.5)
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            draw_organogram(slide, st.session_state.relationships, positions, tree)

            ppt_stream = io.BytesIO()
            prs.save(ppt_stream)
            ppt_stream.seek(0)

            st.success("Apresentação editável gerada com sucesso!")
            st.download_button(
                label="📥 Baixar Apresentação Editável (.pptx)",
                data=ppt_stream,
                file_name="organograma_editavel.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentation.presentation"
            )
else:
    st.warning("Adicione pelo menos uma relação para gerar o organograma.")



